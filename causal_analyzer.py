import os
import json
import threading
import zipfile
import pandas as pd
import magic
from pypdf import PdfReader
import docx
import pptx
from PIL import Image
import pytesseract
from concurrent.futures import ThreadPoolExecutor, as_completed
from datasets import load_dataset
from ollama_models import causal_analyzer, verify_model

# Set the Tesseract-OCR path for image OCR
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
# Max number of concurrent threads (set to 1 for sequential execution)
MAX_WORKERS = 1

# ========== Auxiliary Content Extractors ==========

def load_auxiliary_content(task_id: str, base_folder: str) -> str:
    handlers = [
        (".pdf", _read_pdf_text),
        (".docx", _read_docx_text),
        (".pptx", _read_pptx_text),
        (".xlsx", _read_excel_text),
        (".xls", _read_excel_text),
        (".py", _read_plain_text),
        (".txt", _read_plain_text),
        (".png", _read_image_ocr),
        (".jpg", _read_image_ocr),
        (".jpeg", _read_image_ocr),
        (".zip", _read_zip_listing),
        (".mp3", _note_audio_file),
        (".pdb", _read_plain_or_note),
        (".sls", _read_plain_or_note),
    ]
    for ext, func in handlers:
        path = os.path.join(base_folder, f"{task_id}{ext}")
        if os.path.isfile(path):
            try:
                return func(path)
            except Exception as e:
                return f"[ERROR reading {task_id}{ext}: {e}]"
    return ""

# Extractor implementations (omitted for brevity, same as before)...

def _read_pdf_text(path): return "\n".join((p.extract_text() or "") for p in PdfReader(path).pages).strip()
def _read_docx_text(path): return "\n".join(p.text for p in docx.Document(path).paragraphs).strip()
def _read_pptx_text(path):
    prs, out = pptx.Presentation(path), []
    for slide in prs.slides:
        out.extend([s.text for s in slide.shapes if hasattr(s, "text")])
    return "\n".join(out).strip()
def _read_excel_text(path): return pd.read_excel(path, sheet_name=0).to_string(index=False).strip()
def _read_plain_text(path): return open(path, "r", encoding="utf-8", errors="ignore").read().strip()
def _read_image_ocr(path): return pytesseract.image_to_string(Image.open(path).convert("RGB")).strip()
def _read_zip_listing(path): return "ZIP file contents:\n" + "\n".join(zipfile.ZipFile(path).namelist())
def _note_audio_file(path): return f"[Audio file present: {os.path.basename(path)}]"
def _read_plain_or_note(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read().strip()
    except:
        return f"[Binary or unsupported content: {os.path.basename(path)}]"

# ========== Causal Analysis Logic with Pseudo-Step Generation ==========

def analyze_failure(task_id: str, steps_text: str, aux_text: str, has_steps: bool, max_tokens: int = 2048):
    client = causal_analyzer.client
    # Generate pseudo-steps if no annotator steps
    if not has_steps:
        pseudo_prompt = (
            f"Break down the following question into detailed step-by-step reasoning:\n" 
            f"Question: {steps_text}\nSteps:\n"
        )
        try:
            pseudo_stream = verify_model.client.completion(
                model="ollama/llama3.1:8b",
                provider="ollama",
                api_base="http://localhost:11434",
                api_key="ollama",
                messages=[{"role": "user", "content": pseudo_prompt}],
                max_tokens=max_tokens,
                temperature=0.0,
                top_p=1.0,
                seed=0,
                stream=True
            )
            print(f"--- Task {task_id} Pseudo-steps Generation Start ---")
            pseudo_text = ""
            for chunk in pseudo_stream:
                delta = getattr(chunk.choices[0].delta, "content", "")
                if delta:
                    print(delta, end='', flush=True)
                    pseudo_text += delta
            print(f"--- Task {task_id} Pseudo-steps Generation End ---")
            steps_text = pseudo_text.strip()
        except Exception as e:
            print(f"[WARN] Pseudo-step generation failed for {task_id}: {e}")
    # Build the causal analysis prompt
    prompt = (
        "Below is an Annotator 'Steps' sequence.  \n"
        "NOTE: if these steps were auto-generated, they may contain inaccuracies.\n\n"
        "=== Steps Begin ===\n"
        f"{steps_text.strip()}\n"
        "=== Steps End ===\n\n"
    )
    if aux_text:
        prompt += (
            "=== Auxiliary Document Content Begin ===\n"
            f"{aux_text.strip()}\n"
            "=== Auxiliary Document Content End ===\n\n"
        )
    prompt += (
        "Output format:\n"
        "1. [Top-Level Causal Conclusion]: One-sentence summary of the core inference.\n"
        "2. [Causal Chain]:\n"
        "   a) Step 1 → Why it occurs.\n"
        "   b) Step 2 → Link to Step 1.\n"
        "   ...\n"
        "End with a summary of key causal links."
    )
    # Call causal analyzer with streaming
    response_stream = client.completion(
        model="ollama/nous-hermes2-mixtral:latest",
        provider="ollama",
        api_base="http://localhost:11434",
        api_key="ollama",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=max_tokens,
        temperature=0.0,
        top_p=1.0,
        seed=0,
        stream=True
    )
    print(f"\n=== Task {task_id} Causal Analysis Start ===\n")
    full_text = ""
    for chunk in response_stream:
        delta = getattr(chunk.choices[0].delta, "content", "")
        if delta:
            print(delta, end='', flush=True)
            full_text += delta
    print(f"\n=== Task {task_id} Causal Analysis End ===\n")
    return task_id, full_text

# ========== Main Pipeline ==========

def process_with_gaia_loader():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    gaia_path = os.path.join(base_dir, "GAIA_dataset", "GAIA.py")
    print("📥 Loading GAIA test dataset...")
    dataset = load_dataset(
        path=gaia_path,
        name="2023_all",
        split="test",
        trust_remote_code=True,
    )
    print(f"✅ Successfully loaded {len(dataset)} samples. Starting analysis...")
    tasks = []
    for sample in dataset:
        task_id = sample["task_id"]
        orig_steps = sample.get("Annotator Metadata", {}).get("Steps", "").strip()
        has_steps = bool(orig_steps)
        steps = orig_steps if has_steps else sample.get("Question", "").strip()
        file_path = sample.get("file_path", "")
        aux_text = load_auxiliary_content(task_id, os.path.dirname(file_path)) if file_path else ""
        tasks.append((task_id, steps, aux_text, has_steps))
    if not tasks:
        print("⚠️ No manageable tasks!")
        return
    # Clear output file
    print_lock = threading.Lock()
    open("causal_outputs.jsonl", "w").close()

    def write_result(task_id, result):
        rec = {"task_id": task_id, "causal_analysis": result}
        line = json.dumps(rec, ensure_ascii=False)
        with print_lock:
            with open("causal_outputs.jsonl", "a", encoding="utf-8") as fout:
                fout.write(line + "\n")
                fout.flush()
                os.fsync(fout.fileno())

    # Use thread pool to process tasks in parallel
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = {executor.submit(analyze_failure, *t): t[0] for t in tasks}
        for future in as_completed(futures):
            task_id, result = future.result()
            write_result(task_id, result)

    print("🎉 All tasks have been processed!")

if __name__ == "__main__":
    process_with_gaia_loader()
